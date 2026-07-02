"""Document text extraction utilities supporting multiple formats.

Handles text extraction and page/boundary-chunking for files of type .pdf,
.txt, .md, .docx, .pptx, and .xlsx.
"""

from __future__ import annotations

from pathlib import Path

from logger_utils import setup_logger

logger = setup_logger(__name__)


def _render_page_chunk(page_number: int, content_page_number: int, markdown: str) -> str:
    """Render a single page chunk with machine + human-readable page markers.

    Emits an HTML-comment sentinel (``<!-- page: N -->``) that is invisible in
    rendered markdown but trivially parseable, followed by a ``## Page N``
    heading so agents and the existing ``read_file`` section selector can both
    target an individual page.
    """
    body = markdown.strip()
    return f"<!-- page: {page_number} -->\n## Page {content_page_number}\n\n{body}"


def _infer_page_number_from_boxes(text: str, boxes: list) -> int | None:
    """Attempt to parse a printed page number from header or footer boxes."""
    import re
    for box in boxes:
        if isinstance(box, dict) and box.get("class") in ("page-footer", "page-header"):
            start, end = box.get("pos", (0, 0))
            if start == end:
                continue
            box_text = text[start:end].strip()
            if not box_text:
                continue

            # Match number at beginning (e.g. "**104** BMO Annual Report")
            m_start = re.match(r'^(?:Page\s+)?[*\s\-]*(\d+)[*\s\-]*(?:\s|$)', box_text, re.IGNORECASE)
            if m_start:
                return int(m_start.group(1))

            # Match number at end (e.g. "BMO Annual Report **104**")
            m_end = re.search(r'(?:^|\s)(?:Page\s+)?[*\s\-]*(\d+)[*\s\-]*$', box_text, re.IGNORECASE)
            if m_end:
                return int(m_end.group(1))
    return None


def _extract_pdf_text(file_path: Path) -> str:
    """Extract PDF content as markdown with per-page boundary markers.

    Primary path uses ``pymupdf4llm.to_markdown(page_chunks=True)`` which
    returns a list of per-page dicts (each carrying ``metadata.page``). Each
    page is rendered with a ``<!-- page: N -->`` sentinel + ``## Page N``
    heading so page numbers survive into citations.

    Falls back to the older single-string ``to_markdown`` call (page
    granularity unavailable) and finally to ``pypdf`` per-page extraction.

    Returns:
        str: Extracted markdown content with page markers.
    """
    try:
        import pymupdf4llm

        logger.info("Use PyMuPDF4LLM for PDF markdown extraction (page_chunks=True).")

        markdown_content = pymupdf4llm.to_markdown(
            str(file_path), page_chunks=True, show_progress=False
        )

        if isinstance(markdown_content, list) and markdown_content:
            page_blocks: list[str] = []
            for index, chunk in enumerate(markdown_content, start=1):
                # page_chunks=True returns dicts like {"metadata": {...}, "content": ...}
                if isinstance(chunk, dict):
                    metadata = chunk.get("metadata") or {}
                    page_number = metadata.get("page_number") or metadata.get("page") or index
                    body = chunk.get("text") or chunk.get("content") or chunk.get("markdown") or ""

                    inferred_page = None
                    if "page_boxes" in chunk:
                        inferred_page = _infer_page_number_from_boxes(str(body), chunk["page_boxes"])

                    final_page_num = inferred_page if inferred_page is not None else page_number
                    page_blocks.append(_render_page_chunk(index, int(final_page_num), str(body)))
                else:
                    # Unexpected item shape; render with enumerated page number.
                    page_blocks.append(_render_page_chunk(index, index, str(chunk)))
            return "\n\n".join(page_blocks)

        # Older pymupdf4llm versions or non-chunked mode return a flat string.
        if isinstance(markdown_content, str) and markdown_content.strip():
            logger.warning(
                "pymupdf4llm returned a flat string for %s; "
                "page-level markers unavailable.",
                file_path.name,
            )
            return markdown_content

        # Empty/None result → fall through to pypdf fallback below.
        if markdown_content:
            return str(markdown_content)
    except Exception as e:
        logger.error(f"PyMuPDF4LLM PDF extraction failed: {e}")
        # Fallback to pypdf if markdown extraction fails
        try:
            import pypdf

            logger.info("Falling back to pypdf for PDF text extraction.")
            reader = pypdf.PdfReader(file_path)
            page_texts: list[str] = []
            for index, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    page_texts.append(_render_page_chunk(index, index, text))
            return "\n\n".join(page_texts)
        except Exception as e:
            return f"Error extracting PDF text: {e}"
    return ""


def _extract_text_file(file_path: Path) -> str:
    """Read a plain-text or Markdown file and return its contents.

    Args:
        file_path: Path to a ``.txt`` or ``.md`` file.

    Returns:
        The file contents as a UTF-8 string.
    """
    return file_path.read_text(encoding="utf-8")


def _extract_docx_text(file_path: Path) -> str:
    """Extract DOCX content preserving heading levels as machine sentinels.

    Headings emit an ``<!-- heading: N -->`` sentinel so downstream citation
    capture and navigation can locate claims by heading path.
    """
    from docx import Document

    document = Document(str(file_path))
    lines: list[str] = []

    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        style_name = ""
        try:
            style_name = paragraph.style.name if paragraph.style else ""
        except Exception:
            style_name = ""
        if style_name and style_name.lower().startswith("heading"):
            level = "".join(ch for ch in style_name if ch.isdigit())
            level_int = int(level) if level else 1
            lines.append(f"<!-- heading: {level_int} -->")
            lines.append(text)
        else:
            lines.append(text)

    # Many documents are table-based; include table cells so content is not silently missed.
    table_index = 0
    for table in document.tables:
        table_index += 1
        lines.append(f"<!-- table: {table_index} -->")
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    cell_text = (paragraph.text or "").strip()
                    if cell_text:
                        lines.append(cell_text)

    return "\n".join(lines)


def _extract_pptx_text(file_path: Path) -> str:
    """Extract PPTX content with machine-readable slide sentinels.

    Each slide emits an ``<!-- slide: N -->`` sentinel mirroring the PDF
    ``<!-- page: N -->`` convention, enabling uniform citation capture.
    """
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    slide_sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"<!-- slide: {index} -->", f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes = []
            for shape in slide.notes_slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    notes.append(text.strip())
            notes_text = "\n".join(notes)
        if notes_text:
            parts.append(f"Speaker Notes:\n{notes_text}")
        slide_sections.append("\n".join(parts))

    return "\n\n".join(slide_sections)


def _extract_xlsx_text(file_path: Path) -> str:
    """Extract XLSX content preserving sheet + row locators.

    Each data row carries an ``<!-- sheet: <name>; row: N -->`` sentinel so
    claims can be cited down to a spreadsheet row.
    """
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    sections: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            sheet_lines: list[str] = [f"<!-- sheet: {worksheet.title} -->", f"Sheet: {worksheet.title}"]
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if not values:
                    continue
                sheet_lines.append(f"<!-- sheet: {worksheet.title}; row: {row_index} -->")
                sheet_lines.append(" | ".join(values))
            sections.append(
                "\n".join(sheet_lines) if len(sheet_lines) > 2 else f"Sheet: {worksheet.title}\n(empty sheet)")
    finally:
        workbook.close()

    return "\n\n".join(sections)


def extract_supported_document(file_path: Path) -> str:
    """Extract text content from a supported document file.

    Dispatches to the appropriate extractor based on file extension.
    Supports PDF, plain text, Markdown, DOCX, PPTX, and XLSX.

    Args:
        file_path: Path to the document file.

    Returns:
        Extracted text content, with format-specific sentinels for page
        numbers, slide numbers, headings, or sheet/row coordinates.

    Raises:
        ValueError: If the file extension is not a supported document type.
    """
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_text(file_path)
    if suffix in {".txt", ".md"}:
        return _extract_text_file(file_path)
    if suffix == ".docx":
        return _extract_docx_text(file_path)
    if suffix == ".pptx":
        return _extract_pptx_text(file_path)
    if suffix == ".xlsx":
        return _extract_xlsx_text(file_path)
    raise ValueError(f"Unsupported document type: {suffix}")
